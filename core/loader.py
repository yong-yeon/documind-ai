import os                    # 파일 경로 처리, 파일 존재 확인, 파일 삭제
import re                    # 헤더/푸터 후보 텍스트의 숫자를 정규화 패턴으로 치환하기 위한 정규식
import sys                   # 다른 폴더의 모듈을 import 할 때 경로 추가용
import platform              # OS 판별 (Windows 여부 확인용)
import tempfile              # 업로드 파일을 디스크에 임시 저장하기 위한 임시 파일 생성
import fitz                  # PDF 텍스트 추출 및 페이지 이미지 렌더링 (PyMuPDF)
import pytesseract           # Tesseract OCR 엔진을 파이썬에서 사용하기 위한 래퍼
from PIL import Image        # fitz 픽셀맵 → pytesseract 가 읽을 수 있는 PIL 이미지로 변환

from langchain_core.documents import Document                   # 텍스트 + 메타데이터를 묶는 LangChain 문서 객체
from langchain_text_splitters import RecursiveCharacterTextSplitter  # 긴 텍스트를 청크 단위로 자르는 분할기

# Windows는 Tesseract가 시스템 PATH에 등록되지 않는 경우가 많아 절대 경로를 직접 지정
# Linux(Streamlit Cloud 등)는 packages.txt로 설치된 tesseract-ocr이 PATH에 등록되어 있어 별도 지정 불필요
if platform.system() == "Windows":
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

OCR_THRESHOLD = 20    # fitz 로 추출한 텍스트가 이 글자 수 미만이면 스캔 PDF 로 판단 → OCR 전환
OCR_DPI       = 300   # 300 DPI 로 렌더링해야 한국어 획 복잡도에도 인식률 확보 가능
TESS_CONFIG   = "--psm 6 --oem 3"  # psm 6: 균일 텍스트 블록 모드 / oem 3: LSTM 최신 엔진 사용

HEADER_FOOTER_MIN_RATIO  = 0.3  # 전체 페이지의 이 비율 이상 반복되면 헤더/푸터 후보로 판단
HEADER_FOOTER_MIN_COUNT  = 2    # 비율 계산과 별개로 최소 이 횟수 이상 등장해야 함 (페이지 적은 문서 오탐 방지)
HEADER_FOOTER_MAX_LENGTH = 200  # 이 글자 수 이상인 블록은 본문으로 간주해 헤더/푸터 후보에서 제외

try:
    from config import DEFAULT_CHUNK_SIZE, DEFAULT_CHUNK_OVERLAP  # 앱 정상 실행 시 루트에서 import
except ModuleNotFoundError:
    sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))  # 직접 실행 시 루트 경로 추가
    from config import DEFAULT_CHUNK_SIZE, DEFAULT_CHUNK_OVERLAP  # 경로 추가 후 재시도


def save_uploaded_file(uploaded_file) -> tuple:
    # Streamlit 업로드 객체는 메모리에만 있어서 파일 경로가 없음
    # → 임시 파일로 저장해서 fitz 등이 경로로 접근할 수 있게 함
    suffix = os.path.splitext(uploaded_file.name)[1]  # 원본 파일명에서 확장자만 추출 (예: ".pdf")
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.getvalue())  # 업로드된 바이트 데이터를 임시 파일에 기록
        temp_path = tmp.name                 # 나중에 삭제하기 위해 임시 파일 경로 저장
    return temp_path, uploaded_file.name     # (임시 경로, 원본 파일명) 쌍으로 반환


def _normalize_text(text: str) -> str:
    # 페이지 번호처럼 숫자만 바뀌는 반복 텍스트("Page 1 of 20", "Page 2 of 20")를
    # 같은 패턴("Page # of #")으로 묶어서 비교할 수 있도록 숫자를 # 으로 치환
    return re.sub(r"\d+", "#", text)


def _extract_text_from_blocks(blocks: list, exclude_patterns: set = None) -> str:
    # fitz "blocks" 모드로 얻은 텍스트 블록 목록을 위→아래, 왼→오른쪽 순서로 정렬
    # exclude_patterns 에 포함된 정규화 패턴(반복 헤더/푸터로 판단된 블록)은 제외하고 본문만 저장
    exclude_patterns = exclude_patterns or set()

    text_blocks = []
    for b in sorted(blocks, key=lambda b: (round(b[1] / 10), b[0])):  # y좌표→x좌표 순 정렬 (읽기 순서 보정)
        text = b[4].strip()
        if b[6] != 0 or not text:            # 이미지 블록(type≠0) 이거나 빈 블록이면 건너뜀
            continue
        if _normalize_text(text) in exclude_patterns:  # 반복 헤더/푸터 패턴과 일치하면 제외
            continue
        text_blocks.append(b[4])             # 본문 블록 텍스트 수집
    return "\n".join(text_blocks)            # 블록들을 줄바꿈으로 합쳐서 반환


def _scan_pages(doc, min_ratio: float = HEADER_FOOTER_MIN_RATIO,
                 min_count: int = HEADER_FOOTER_MIN_COUNT, max_length: int = HEADER_FOOTER_MAX_LENGTH) -> tuple:
    # PDF 전체 페이지를 한 번만 스캔해서
    #   1) 페이지별 블록 목록을 캐싱 (load_pdf 에서 재사용 → 블록 추출 중복 방지)
    #   2) 페이지의 min_ratio 이상 & 최소 min_count 회 이상 반복되는 짧은 텍스트 패턴을 헤더/푸터로 판단
    total_pages = doc.page_count
    page_blocks_list = []                    # 페이지별 원본 블록 목록 (인덱스 = page_num)
    pattern_page_count = {}                  # 정규화 패턴별로 등장한 "페이지 수" 집계 (한 페이지 내 중복은 1회로 계산)

    for page in doc:
        blocks = page.get_text("blocks")     # 텍스트 블록 목록 추출 (x0,y0,x1,y1,text,no,type)
        page_blocks_list.append(blocks)

        patterns_on_page = set()
        for b in blocks:
            if b[6] != 0:                    # 이미지 블록 제외
                continue
            text = b[4].strip()
            if not text or len(text) >= max_length:  # 빈 텍스트 제외, 200자 이상 제외 (짧은 텍스트만 대상)
                continue
            patterns_on_page.add(_normalize_text(text))
        for pattern in patterns_on_page:
            pattern_page_count[pattern] = pattern_page_count.get(pattern, 0) + 1

    min_pages = max(min_count, total_pages * min_ratio)  # 비율 조건과 최소 횟수 조건 중 더 엄격한 쪽 적용
    repeated_patterns = {pattern for pattern, count in pattern_page_count.items() if count >= min_pages}
    return page_blocks_list, repeated_patterns


def _get_page_label(page, page_num: int) -> str:
    # PDF 에 실제 인쇄된 페이지 번호 (예: "i", "1", "A-1") 를 반환
    # 레이블이 없으면 물리적 순서 기반 번호 (1-index) 로 대체
    try:
        label = page.get_label()             # PDF 내부에 저장된 페이지 레이블 읽기
        return label.strip() if label.strip() else str(page_num + 1)  # 있으면 레이블, 없으면 순서+1
    except Exception:
        return str(page_num + 1)             # 오류 시 물리적 순서로 대체


def load_pdf(file_path: str, original_name: str) -> list:
    doc = fitz.open(file_path)              # PDF 파일 열기
    page_blocks_list, repeated_patterns = _scan_pages(doc)  # 블록 캐싱 + 반복 헤더/푸터 패턴 탐지 (전체 페이지 1회 스캔)
    docs = []                               # 결과 Document 담을 빈 리스트
    for page_num, page in enumerate(doc):  # 페이지 번호와 페이지 객체를 함께 순회
        blocks   = page_blocks_list[page_num]                # 스캔 단계에서 캐싱해둔 블록 재사용 (재추출 방지)
        raw_text = _extract_text_from_blocks(blocks)          # 전체 텍스트 추출 (OCR 판단용, 반복 블록 제외 전)

        if len(raw_text.strip()) >= OCR_THRESHOLD:  # 20글자 이상 → 일반 PDF
            clean_text = _extract_text_from_blocks(blocks, exclude_patterns=repeated_patterns)  # 반복 헤더/푸터 제외한 본문
            page_label = _get_page_label(page, page_num)                       # 실제 인쇄 페이지 번호
            docs.append(Document(
                page_content=clean_text,
                metadata={"source": original_name, "page": page_num, "page_label": page_label}
            ))
        else:                                        # 20글자 미만 → 스캔 PDF → OCR 전환
            try:
                pix      = page.get_pixmap(dpi=OCR_DPI)                        # 300 DPI 로 페이지 이미지 렌더링
                img      = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)  # PIL 이미지로 변환
                ocr_text = pytesseract.image_to_string(img, lang="kor+eng", config=TESS_CONFIG)  # OCR 실행
            except Exception:
                ocr_text = ""                        # OCR 실패 시 빈 문자열로 조용히 건너뜀
            if ocr_text.strip():                     # OCR 결과가 있을 때만 저장
                page_label = _get_page_label(page, page_num)
                docs.append(Document(
                    page_content=ocr_text,
                    metadata={"source": original_name, "page": page_num, "page_label": page_label}
                ))
    doc.close()                                      # PDF 파일 핸들 닫기 (메모리 해제)
    return docs


def load_docx(file_path: str, original_name: str) -> list:
    from docx import Document as DocxDocument        # 함수 내 import 로 불필요한 의존성 최소화
    docx      = DocxDocument(file_path)              # DOCX 파일 열기
    full_text = []

    for para in docx.paragraphs:                     # 문서의 단락(문단) 하나씩 순회
        if para.text.strip():                        # 빈 단락 제외
            full_text.append(para.text)

    for table in docx.tables:                        # 문서의 표 순회
        for row in table.rows:                       # 표의 행 순회
            row_text = " | ".join(                   # 셀 내용을 " | " 구분자로 연결
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                full_text.append(row_text)

    docs = []
    if full_text:
        docs.append(Document(
            page_content="\n".join(full_text),       # 전체 내용을 하나의 Document 로 합침
            metadata={"source": original_name, "page": 0}  # DOCX 는 페이지 구분 없어서 0 고정
        ))
    return docs


def load_xlsx(file_path: str, original_name: str) -> list:
    import openpyxl                                  # 함수 내 import 로 의존성 최소화
    wb   = openpyxl.load_workbook(file_path, data_only=True)  # 수식 대신 계산 결과값으로 읽기
    docs = []

    for sheet_num, sheet_name in enumerate(wb.sheetnames):  # 시트 번호와 이름 함께 순회
        ws        = wb[sheet_name]                   # 현재 시트 선택
        rows_text = []
        for row in ws.iter_rows(values_only=True):   # 행별로 값만 추출
            row_data = [str(cell) for cell in row if cell is not None]  # None 셀 제외 후 문자열 변환
            if row_data:
                rows_text.append(" | ".join(row_data))  # 셀을 " | " 구분자로 연결

        if rows_text:
            docs.append(Document(
                page_content=f"[시트: {sheet_name}]\n" + "\n".join(rows_text),  # 시트명 헤더 포함
                metadata={"source": original_name, "page": sheet_num}           # 시트 번호를 페이지로 사용
            ))
    return docs


def load_pptx(file_path: str, original_name: str) -> list:
    from pptx import Presentation                    # 함수 내 import 로 의존성 최소화
    prs  = Presentation(file_path)                   # PPTX 파일 열기
    docs = []

    for slide_num, slide in enumerate(prs.slides):  # 슬라이드 번호와 슬라이드 객체 함께 순회
        slide_text = []
        for shape in slide.shapes:                   # 슬라이드 안의 도형 하나씩 순회
            if shape.has_text_frame:                 # 텍스트 박스 도형인 경우
                for para in shape.text_frame.paragraphs:  # 텍스트 박스 안의 단락 순회
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
            if shape.has_table:                      # 표 도형인 경우
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        slide_text.append(row_text)

        if slide_text:
            docs.append(Document(
                page_content=f"[슬라이드 {slide_num + 1}]\n" + "\n".join(slide_text),  # 슬라이드 번호 헤더 포함
                metadata={"source": original_name, "page": slide_num}
            ))
    return docs


def load_file(file_input, chunk_size: int = DEFAULT_CHUNK_SIZE, chunk_overlap: int = DEFAULT_CHUNK_OVERLAP):
    # 로컬 파일 경로(str) 또는 Streamlit 업로드 객체를 받아서 청크 리스트로 반환
    temp_path = None
    if isinstance(file_input, str):                  # 로컬 파일 경로로 전달된 경우
        if not os.path.exists(file_input):
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_input}")
        file_path     = file_input
        original_name = os.path.basename(file_input) # 전체 경로에서 파일명만 추출
    else:                                            # Streamlit 업로드 객체로 전달된 경우
        temp_path, original_name = save_uploaded_file(file_input)  # 임시 파일로 저장 후 경로 반환
        file_path = temp_path

    try:
        ext = os.path.splitext(file_path)[1].lower() # 확장자 추출 후 소문자 통일 (예: ".PDF" → ".pdf")
        if ext == ".pdf":
            docs = load_pdf(file_path, original_name)
        elif ext == ".docx":
            docs = load_docx(file_path, original_name)
        elif ext == ".xlsx":
            docs = load_xlsx(file_path, original_name)
        elif ext in [".pptx", ".ppt"]:
            docs = load_pptx(file_path, original_name)
        else:
            raise ValueError(f"지원하지 않는 파일 형식입니다: {ext}")

        if not docs:
            raise ValueError("문서에서 텍스트를 추출할 수 없습니다.")

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,               # 청크 최대 글자 수 (기본 1000자)
            chunk_overlap=chunk_overlap,         # 청크 간 겹치는 글자 수 (기본 150자, 문맥 연결용)
            separators=["\n\n", "\n", ". ", "다. ", "요. ", "! ", "? ", " ", ""]  # 자연스러운 경계 우선 분할
        )
        split_docs = splitter.split_documents(docs)  # 추출된 Document 리스트를 청크로 분할

    finally:
        if temp_path and os.path.exists(temp_path):  # 임시 파일이 있으면
            try:
                os.remove(temp_path)                 # 성공/실패 관계없이 항상 임시 파일 삭제
            except (PermissionError, OSError):
                pass                                 # Windows 에서 파일 잠금 오류 시 조용히 무시

    return split_docs                                # 청크로 나뉜 Document 리스트 반환


def get_doc_info(split_docs) -> dict:
    # 청크 목록에서 문서 요약 정보를 추출해 반환 (디버깅/테스트용)
    if not split_docs:
        return {"총 조각 수": 0, "총 페이지 수": 0, "출처 파일": "없음", "첫 번째 조각": ""}

    pages = [doc.metadata.get("page", 0) for doc in split_docs]  # 모든 청크의 페이지 번호 수집
    return {
        "총 조각 수":   len(split_docs),
        "총 페이지 수": max(pages) + 1,              # 0-index 이므로 +1 해서 실제 페이지 수 계산
        "출처 파일":    split_docs[0].metadata.get("source", "알 수 없음"),
        "첫 번째 조각": split_docs[0].page_content[:100],  # 첫 번째 청크 미리보기 100자
    }